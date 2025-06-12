from pathlib import Path

from icecream import ic
from pptx import Presentation

TARGET_MD = "06_lec.md"


def prepare_slide_objects():
    # テンプレートファイルのパス
    TEMPLATE_PPTM = current_dir / "slide_template.pptm"

    # プレゼンテーションをテンプレートから読み込む
    return Presentation(TEMPLATE_PPTM)


def convert_slide_contents(read_lines):
    slide_dict = {}
    slide_title_list = []
    for line in read_lines:
        # スライドのタイトルと内容を取得
        if line.startswith("- "):
            slide_title = line[2:].strip()
            slide_dict[slide_title] = []
            slide_title_list.append(slide_title)
            continue
        slide_dict[slide_title].append(line.strip("- |\n"))
    # ic(slide_dict)

    return slide_dict, slide_title_list


def generate_slides(base_presentation):
    slide_dict, slide_title_list = convert_slide_contents(lines)

    for slide_title_key in slide_title_list:
        slide_content = slide_dict[slide_title_key]
        # スライドのタイトルとコンテンツ
        slide_title = slide_title_key
        # slide_content = ["箇条書き項目1", "箇条書き項目2", "箇条書き項目3"]

        # 新しいスライドを追加
        slide_layout = base_presentation.slide_layouts[
            2
        ]  # タイトルとコンテンツのレイアウト
        slide = base_presentation.slides.add_slide(slide_layout)

        # タイトルを設定
        title = slide.shapes.title
        title.text = slide_title

        # コンテンツを設定
        content = slide.placeholders[1]
        paragraph_tf = content.text_frame
        for item in slide_content:
            p = paragraph_tf.add_paragraph()  # 新しい段落を追加
            p.text = item


if __name__ == "__main__":
    current_dir = Path(__file__).resolve().parents[0]

    target_outline = current_dir / TARGET_MD
    with open(target_outline, "r", encoding="utf-8") as f:
        lines = f.readlines()

    presentation = prepare_slide_objects()
    generate_slides(presentation)

    # ファイルを保存
    output_path = str(target_outline).replace(".md", ".pptm")
    presentation.save(output_path)
    print(f"スライドが保存されました: {output_path}")
